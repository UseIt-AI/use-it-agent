"""
PPT Action Executor - 结构化 Action 解析与执行

统一调度入口，将大模型输出的 JSON actions 路由到对应的处理函数。

支持的 Action 类型:

[幻灯片管理]
  add_slide          - 添加幻灯片
  delete_slide       - 删除幻灯片
  duplicate_slide    - 复制幻灯片
  move_slide         - 移动幻灯片

[感知]
  get_current_context - 获取当前幻灯片完整上下文（占位符、元素、样式）

[SVG 渲染]
  render_ppt_layout  - render_svg_layout 的别名

[原生复杂对象]
  insert_native_chart - 插入原生图表
  insert_native_table - 插入原生表格
  insert_media        - 插入图片/媒体

[状态编辑]
  update_element     - 按 handle_id 修改已有元素属性
                       支持: 位置/尺寸、fill_color/fill_gradient、
                       line_color/line_weight、shadow、text/font 属性
  delete_element     - 按 handle_id 删除元素
  group_elements     - 将多个元素组合为一个 Group
  ungroup_elements   - 解散一个 Group

[对齐 / 分布]
  align_elements       - 对齐多个元素
  distribute_elements  - 均匀分布多个元素

[动画 / 切换]
  add_shape_animation     - 添加形状动画
  clear_slide_animations  - 清除幻灯片动画
  set_slide_transition    - 设置幻灯片切换效果
"""

import logging
import os
import tempfile
from typing import Any, Dict, List, Optional

import math

from .constants import (
    MSO_GRADIENT_FROM_CENTER,
    MSO_GRADIENT_HORIZONTAL,
    PLACEHOLDER_TYPES,
    SHAPE_TYPE_NAMES,
    SHAPE_TYPES,
    SLIDE_LAYOUTS,
    TEXT_ALIGN,
    Z_ORDER_COMMANDS,
    color_int_to_hex,
    find_shape_by_handle,
    find_shape_by_index,
    parse_color,
    resolve_slide,
    resolve_slide_with_app,
    resolve_table_style,
)
from .animation_constants import (
    resolve_mso_anim_effect,
    resolve_mso_anim_text_unit_effect,
    resolve_mso_anim_trigger,
    resolve_pp_entry_effect,
)
from .macro_executor import MacroExecutor
from .renderer import SlideRenderer
from .renderer.linter import LayoutLinter

logger = logging.getLogger(__name__)


class PPTActionExecutor:
    """
    结构化 Action 执行器

    接收 actions 列表，按顺序执行，返回每个 action 的结果。
    """

    def __init__(self):
        self.svg_renderer = SlideRenderer()
        self.macro_executor = MacroExecutor()

    def execute_actions(
        self, app, pres, actions: List[Dict[str, Any]]
    ) -> Dict[str, Any]:
        """
        批量执行 actions。

        Args:
            app: COM PowerPoint.Application object
            pres: COM Presentation object
            actions: action 列表

        Returns:
            {
                "success": bool,       # 全部成功为 True
                "results": [...],      # 每个 action 的结果
                "error": str|None,     # 第一个失败的错误信息
            }
        """
        results = []
        overall_success = True
        first_error = None

        for i, action in enumerate(actions):
            action_name = action.get("action")
            if not action_name:
                result = {"success": False, "error": "Missing 'action' field"}
            else:
                handler = self._get_handler(action_name)
                if handler is None:
                    result = {
                        "success": False,
                        "error": f"Unknown action: {action_name}",
                    }
                else:
                    try:
                        result = handler(app, pres, action)
                    except Exception as e:
                        logger.error(
                            f"[ActionExecutor] Action '{action_name}' (#{i}) failed: {e}",
                            exc_info=True,
                        )
                        result = {"success": False, "error": str(e)}

            results.append({"action": action_name, "index": i, **result})

            if not result.get("success", False):
                overall_success = False
                if first_error is None:
                    first_error = f"Action #{i} ({action_name}): {result.get('error', 'unknown')}"
                # Don't break on error — continue executing remaining actions
                # (the LLM might have independent actions after a failed one)

        return {
            "success": overall_success,
            "results": results,
            "error": first_error,
        }

    def _get_handler(self, action_name: str):
        """Map action name to handler method."""
        return {
            "add_slide": self._handle_add_slide,
            "delete_slide": self._handle_delete_slide,
            "duplicate_slide": self._handle_duplicate_slide,
            "move_slide": self._handle_move_slide,
            "goto_slide": self._handle_goto_slide,
            "render_ppt_layout": self._handle_render_ppt_layout,
            "lint_layout": self._handle_lint_layout,
            "insert_native_chart": self._handle_insert_native_chart,
            "insert_native_table": self._handle_insert_native_table,
            "insert_media": self._handle_insert_media,
            "update_element": self._handle_update_element,
            "delete_element": self._handle_delete_element,
            "group_elements": self._handle_group_elements,
            "ungroup_elements": self._handle_ungroup_elements,
            "align_elements": self._handle_align_elements,
            "distribute_elements": self._handle_distribute_elements,
            "reorder_elements": self._handle_reorder_elements,
            "add_shape_animation": self._handle_add_shape_animation,
            "clear_slide_animations": self._handle_clear_slide_animations,
            "set_slide_transition": self._handle_set_slide_transition,
        }.get(action_name)

    # ==================== 幻灯片管理 ====================

    def _handle_add_slide(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        添加幻灯片。

        Params:
            layout: str|int  - 版式名称（见 SLIDE_LAYOUTS）或 ppLayoutType 整数
            index: int       - 插入位置（可选，默认追加到末尾）
        """
        layout = params.get("layout", "blank")
        index = params.get("index")

        if isinstance(layout, str):
            layout_id = SLIDE_LAYOUTS.get(layout.lower())
            if layout_id is None:
                return {
                    "success": False,
                    "error": f"Unknown layout: {layout}. Available: {list(SLIDE_LAYOUTS.keys())}",
                }
        else:
            layout_id = int(layout)

        if index is None:
            index = pres.Slides.Count + 1

        slide = pres.Slides.Add(index, layout_id)

        try:
            app.ActiveWindow.View.GotoSlide(slide.SlideIndex)
        except Exception:
            pass

        logger.info(f"[ActionExecutor] Added slide at index {slide.SlideIndex}, layout={layout}")

        return {
            "success": True,
            "slide_index": slide.SlideIndex,
        }

    def _handle_goto_slide(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        跳转到指定幻灯片。

        Params:
            slide: int|"first"|"last"|"next"|"prev" - 目标幻灯片
        """
        target = params.get("slide")
        if target is None:
            return {"success": False, "error": "Missing 'slide' parameter"}

        total = pres.Slides.Count
        if total == 0:
            return {"success": False, "error": "No slides in presentation"}

        current = app.ActiveWindow.View.Slide.SlideIndex

        if isinstance(target, str):
            target_map = {
                "first": 1,
                "last": total,
                "next": min(current + 1, total),
                "prev": max(current - 1, 1),
            }
            idx = target_map.get(target.lower())
            if idx is None:
                return {"success": False, "error": f"Invalid slide target: {target}. Use int, first, last, next, prev"}
        else:
            idx = int(target)

        if idx < 1 or idx > total:
            return {"success": False, "error": f"Slide index {idx} out of range (1-{total})"}

        app.ActiveWindow.View.GotoSlide(idx)
        logger.info(f"[ActionExecutor] Navigated to slide {idx}")

        return {"success": True, "slide_index": idx}

    def _handle_delete_slide(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        删除幻灯片。

        Params:
            slide: int - 要删除的幻灯片索引（1-based）
        """
        slide_ref = params.get("slide")
        if slide_ref is None:
            return {"success": False, "error": "Missing 'slide' parameter"}

        slide = resolve_slide_with_app(app, pres, slide_ref)
        slide_idx = slide.SlideIndex
        slide.Delete()
        logger.info(f"[ActionExecutor] Deleted slide {slide_idx}")

        return {"success": True, "deleted_index": slide_idx}

    def _handle_duplicate_slide(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        复制幻灯片。

        Params:
            slide: int - 要复制的幻灯片索引
        """
        slide_ref = params.get("slide")
        if slide_ref is None:
            return {"success": False, "error": "Missing 'slide' parameter"}

        slide = resolve_slide(pres, slide_ref)
        new_slide = slide.Duplicate()
        new_index = new_slide(1).SlideIndex  # Duplicate returns SlideRange
        logger.info(f"[ActionExecutor] Duplicated slide {slide_ref} → {new_index}")

        return {"success": True, "new_slide_index": new_index}

    def _handle_move_slide(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        移动幻灯片。

        Params:
            slide: int    - 要移动的幻灯片索引
            to_index: int - 目标位置
        """
        slide_ref = params.get("slide")
        to_index = params.get("to_index")
        if slide_ref is None or to_index is None:
            return {"success": False, "error": "Missing 'slide' or 'to_index'"}

        slide = resolve_slide(pres, slide_ref)
        slide.MoveTo(to_index)
        logger.info(f"[ActionExecutor] Moved slide {slide_ref} → {to_index}")

        return {"success": True}

    # ==================== SVG 渲染 ====================

    def _handle_render_ppt_layout(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        将 SVG 渲染为 PPT 形状。

        Params:
            slide: int|"current"  - 目标幻灯片
            svg: str              - SVG 字符串
            clear_slide: bool     - 是否先清空幻灯片（默认 False）
            render_mode: str      - "create" | "supplement" | "patch"
                create:     clear_slide=True 全量重建
                supplement: clear_slide=False 追加新 shape（默认）
                patch:      就地编辑已有元素（按 data-handle-id 匹配）
        """
        slide_ref = params.get("slide", "current")
        svg_string = params.get("svg")
        clear_slide = params.get("clear_slide", False)
        render_mode = params.get("render_mode")
        render_strategy = params.get("render_strategy")
        layers = params.get("layers")
        palette = params.get("palette")
        patch_scope = params.get("patch_scope")

        if not svg_string:
            return {"success": False, "error": "Missing 'svg' parameter"}

        # Backward compatible: infer render_mode from clear_slide
        if render_mode is None:
            render_mode = "create" if clear_slide else "supplement"

        slide = resolve_slide_with_app(app, pres, slide_ref)
        slide_width = pres.PageSetup.SlideWidth
        slide_height = pres.PageSetup.SlideHeight

        if render_mode == "create":
            for i in range(slide.Shapes.Count, 0, -1):
                try:
                    slide.Shapes(i).Delete()
                except Exception:
                    pass
            if render_strategy == "layered":
                return self.svg_renderer.render_layered(
                    slide, svg_string, slide_width, slide_height, layers, palette
                )
            return self.svg_renderer.render(slide, svg_string, slide_width, slide_height)

        if render_mode == "patch":
            if render_strategy == "layered" or (
                isinstance(patch_scope, dict) and patch_scope.get("type") == "layer"
            ):
                return self.svg_renderer.patch_layered(
                    slide, svg_string, slide_width, slide_height,
                    layers, patch_scope, palette,
                )
            return self.svg_renderer.patch(slide, svg_string, slide_width, slide_height)

        # supplement (default)
        if render_strategy == "layered":
            return self.svg_renderer.render_layered(
                slide, svg_string, slide_width, slide_height, layers, palette
            )
        return self.svg_renderer.render(slide, svg_string, slide_width, slide_height)

    # ==================== 布局校验 ====================

    def _handle_lint_layout(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """Validate the slide against its declared layer blueprint.

        Params:
            slide: int|"current" - 目标幻灯片
            blueprint: list[dict] (可选) - 覆盖 slide 上保存的 blueprint
                                            （未指定时从 slide.Tags 读取）

        Returns:
            success, issues, shape_count, layer_count, blueprint
        """
        slide_ref = params.get("slide", "current")
        slide = resolve_slide_with_app(app, pres, slide_ref)
        slide_w = pres.PageSetup.SlideWidth
        slide_h = pres.PageSetup.SlideHeight

        blueprint = params.get("blueprint")
        if blueprint is None:
            blueprint = SlideRenderer._load_blueprint(slide)

        shapes_info = self._collect_shapes_for_lint(slide)
        issues = LayoutLinter.lint(shapes_info, slide_w, slide_h, blueprint)
        return {
            "success": True,
            "issues": issues,
            "shape_count": len(shapes_info),
            "layer_count": len({s["layer_id"] for s in shapes_info if s.get("layer_id")}),
            "blueprint": blueprint or [],
        }

    @staticmethod
    def _collect_shapes_for_lint(slide) -> List[Dict[str, Any]]:
        """Read every shape on the slide as a minimal lint-ready dict.

        ``layer_id`` resolution priority: shape.Tags("useit_layer_id") →
        prefix of shape.Name before the first dot. Solid fill / line colors
        are extracted (as ``#RRGGBB``) so the palette-drift rule can run.
        """
        out: List[Dict[str, Any]] = []
        for i in range(1, slide.Shapes.Count + 1):
            try:
                shape = slide.Shapes(i)
            except Exception:
                continue
            try:
                name = shape.Name or ""
            except Exception:
                name = ""
            layer_id = None
            try:
                tag_val = shape.Tags("useit_layer_id")
                layer_id = (tag_val or None) or layer_id
            except Exception:
                pass
            if not layer_id and "." in name:
                prefix = name.split(".", 1)[0].strip()
                layer_id = prefix or None
            try:
                bounds = {
                    "x": round(shape.Left, 2),
                    "y": round(shape.Top, 2),
                    "w": round(shape.Width, 2),
                    "h": round(shape.Height, 2),
                }
            except Exception:
                bounds = None
            fill_color = None
            try:
                fill = shape.Fill
                if fill.Type == 1:  # msoFillSolid
                    fill_color = color_int_to_hex(fill.ForeColor.RGB)
                elif fill.Type == 5:  # msoFillBackground
                    fill_color = "none"
            except Exception:
                pass
            line_color = None
            try:
                line = shape.Line
                if line.Visible:
                    line_color = color_int_to_hex(line.ForeColor.RGB)
            except Exception:
                pass
            out.append({
                "index": i,
                "handle_id": name or None,
                "layer_id": layer_id,
                "bounds": bounds,
                "fill_color": fill_color,
                "line_color": line_color,
            })
        return out

    # ==================== 原生复杂对象 ====================

    def _handle_insert_native_chart(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        插入原生图表。

        Params:
            slide: int|"current"
            chart_type: str
            bounding_box: {"x","y","w","h"}
            data: [[...], [...]]  (二维数组)
            title: str (可选)
            style: int (可选)
            handle_id: str (可选)
        """
        slide_ref = params.get("slide", "current")
        slide = resolve_slide_with_app(app, pres, slide_ref)

        return self.macro_executor.insert_chart(
            slide=slide,
            chart_type=params.get("chart_type", "column_clustered"),
            bounds=params.get("bounding_box", {}),
            data_matrix=params.get("data", []),
            title=params.get("title"),
            style=params.get("style"),
            handle_id=params.get("handle_id"),
        )

    def _handle_insert_native_table(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        插入原生表格。

        Params:
            slide: int|"current"
            bounding_box: {"x","y","w","h"}
            data: [[...], [...]]
            style: str — 样式名称或 GUID（默认 "no_style"）
            first_row_header: bool (默认 True)
            first_col: bool (默认 False)
            banded_rows: bool (默认 True)
            banded_cols: bool (默认 False)
            font_size: float (可选，全表默认字号)
            cell_format: [{row,col,fill_color,font_color,font_bold,...}]
            handle_id: str (可选)
        """
        slide_ref = params.get("slide", "current")
        slide = resolve_slide_with_app(app, pres, slide_ref)

        return self.macro_executor.insert_table(
            slide=slide,
            bounds=params.get("bounding_box", {}),
            data_matrix=params.get("data", []),
            style=params.get("style"),
            first_row_header=params.get("first_row_header", True),
            first_col=params.get("first_col", False),
            banded_rows=params.get("banded_rows", True),
            banded_cols=params.get("banded_cols", False),
            handle_id=params.get("handle_id"),
            cell_format=params.get("cell_format"),
            font_size=params.get("font_size"),
        )

    def _handle_insert_media(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        插入媒体文件（图片/视频）。

        Params:
            slide: int|"current"
            media_path: str     - 文件绝对路径
            bounding_box: {"x","y","w","h"}
            handle_id: str (可选)
            preserve_aspect_ratio: str (可选) — 与 SVG preserveAspectRatio 相同；默认 xMidYMid meet
        """
        slide_ref = params.get("slide", "current")
        slide = resolve_slide_with_app(app, pres, slide_ref)
        media_path = params.get("media_path")

        if not media_path:
            return {"success": False, "error": "Missing 'media_path' parameter"}

        return self.macro_executor.insert_media(
            slide=slide,
            media_path=media_path,
            bounds=params.get("bounding_box", {}),
            handle_id=params.get("handle_id"),
            preserve_aspect_ratio=params.get("preserve_aspect_ratio"),
        )

    # ==================== 状态编辑 ====================

    def _handle_update_element(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        按 handle_id 修改已有元素属性。

        Params:
            slide: int|"current"
            handle_id: str
            properties: dict  — 支持的 key 见 ``_apply_properties``
        """
        slide_ref = params.get("slide", "current")
        handle_id = params.get("handle_id")
        properties = params.get("properties", {})

        if not handle_id:
            return {"success": False, "error": "Missing 'handle_id'"}
        if not properties:
            return {"success": False, "error": "Missing 'properties'"}

        slide = resolve_slide_with_app(app, pres, slide_ref)
        shape = find_shape_by_handle(slide, handle_id)

        if shape is None:
            return {
                "success": False,
                "error": f"Shape not found: '{handle_id}' on slide {slide.SlideIndex}",
            }

        try:
            self._apply_properties(shape, properties)
            logger.info(
                f"[ActionExecutor] Updated element '{handle_id}': "
                f"{list(properties.keys())}"
            )
            return {"success": True, "handle_id": handle_id}
        except Exception as e:
            return {"success": False, "error": str(e)}

    def _handle_delete_element(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        按 handle_id 删除元素。

        Params:
            slide: int|"current"
            handle_id: str
        """
        slide_ref = params.get("slide", "current")
        handle_id = params.get("handle_id")

        if not handle_id:
            return {"success": False, "error": "Missing 'handle_id'"}

        slide = resolve_slide_with_app(app, pres, slide_ref)
        shape = find_shape_by_handle(slide, handle_id)

        if shape is None:
            return {
                "success": False,
                "error": f"Shape not found: '{handle_id}'",
            }

        shape.Delete()
        logger.info(f"[ActionExecutor] Deleted element '{handle_id}'")
        return {"success": True, "deleted": handle_id}

    def _handle_group_elements(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        将多个元素组合为一个 Group。

        Params:
            slide: int|"current"
            handle_ids: list[str]   - 要组合的 Shape.Name 列表（至少 2 个）
            group_name: str         - 组合后的 Group 名称（可选）
        """
        slide_ref = params.get("slide", "current")
        handle_ids = params.get("handle_ids", [])
        group_name = params.get("group_name")

        if len(handle_ids) < 2:
            return {"success": False, "error": "Need at least 2 handle_ids to group"}

        slide = resolve_slide_with_app(app, pres, slide_ref)

        indices = []
        missing = []
        for hid in handle_ids:
            found = False
            for i in range(1, slide.Shapes.Count + 1):
                try:
                    if slide.Shapes(i).Name == hid:
                        indices.append(i)
                        found = True
                        break
                except Exception:
                    continue
            if not found:
                missing.append(hid)

        if missing:
            return {
                "success": False,
                "error": f"Shapes not found: {missing}",
            }

        try:
            shape_range = slide.Shapes.Range(indices)
            group = shape_range.Group()

            if group_name:
                group.Name = group_name

            result_name = group.Name
            logger.info(
                f"[ActionExecutor] Grouped {len(handle_ids)} elements → '{result_name}'"
            )
            return {"success": True, "handle_id": result_name}
        except Exception as e:
            return {"success": False, "error": str(e)}

    def _handle_ungroup_elements(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        解散一个 Group，释放其中的子元素。

        Params:
            slide: int|"current"
            handle_id: str  - Group 的 Shape.Name
        """
        slide_ref = params.get("slide", "current")
        handle_id = params.get("handle_id")

        if not handle_id:
            return {"success": False, "error": "Missing 'handle_id'"}

        slide = resolve_slide_with_app(app, pres, slide_ref)
        shape = find_shape_by_handle(slide, handle_id)

        if shape is None:
            return {
                "success": False,
                "error": f"Shape not found: '{handle_id}'",
            }

        try:
            ungrouped = shape.Ungroup()
            child_names = []
            for i in range(1, ungrouped.Count + 1):
                child_names.append(ungrouped(i).Name)
            logger.info(
                f"[ActionExecutor] Ungrouped '{handle_id}' → {len(child_names)} elements"
            )
            return {"success": True, "children": child_names}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ==================== Align / Distribute ====================

    def _handle_align_elements(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        Align elements using manual coordinate calculation.

        Supports 4 modes (determined by parameter priority):
          1. targets + references[]  → 1:1 paired alignment
          2. targets + reference (shape name) → align all targets to that shape
          3. targets + reference="slide"      → align all targets to slide
          4. targets only (no reference)      → mutual alignment

        Params:
            slide: int|"current"
            targets: list[str]          - shapes to move
            reference: str|None         - single reference shape name or "slide"
            references: list[str]|None  - array of reference shapes (1:1 with targets)
            horizontal: str|None        - "left"/"center"/"right"
            vertical: str|None          - "top"/"middle"/"bottom"

            # Legacy params (backward compat)
            handle_ids: list[str]       - alias for targets
            alignment: str              - mapped to horizontal or vertical
            relative_to: str            - "slide" → reference="slide"
        """
        slide_ref = params.get("slide", "current")
        slide = resolve_slide_with_app(app, pres, slide_ref)

        targets = params.get("targets") or params.get("handle_ids", [])
        reference = params.get("reference")
        references = params.get("references")
        horizontal = params.get("horizontal")
        vertical = params.get("vertical")

        if horizontal is None and vertical is None:
            alignment = params.get("alignment")
            if alignment:
                if alignment in ("left", "center", "right"):
                    horizontal = alignment
                elif alignment in ("top", "middle", "bottom"):
                    vertical = alignment
            rel = params.get("relative_to")
            if rel == "slide" and reference is None:
                reference = "slide"

        if not targets:
            return {"success": False, "error": "No targets specified"}
        if horizontal is None and vertical is None:
            return {"success": True, "aligned_count": 0, "skipped": []}

        VALID_H = ("left", "center", "right")
        VALID_V = ("top", "middle", "bottom")
        if horizontal is not None and horizontal not in VALID_H:
            return {"success": False, "error": f"Invalid horizontal: {horizontal}. Use: {VALID_H}"}
        if vertical is not None and vertical not in VALID_V:
            return {"success": False, "error": f"Invalid vertical: {vertical}. Use: {VALID_V}"}

        def _find_shape(name: str):
            for i in range(1, slide.Shapes.Count + 1):
                try:
                    if slide.Shapes(i).Name == name:
                        return slide.Shapes(i)
                except Exception:
                    continue
            return None

        def _bounds(shape):
            return {
                "left": float(shape.Left), "top": float(shape.Top),
                "width": float(shape.Width), "height": float(shape.Height),
            }

        def _align_one(target_shape, ref_b):
            if horizontal == "left":
                target_shape.Left = ref_b["left"]
            elif horizontal == "center":
                target_shape.Left = ref_b["left"] + ref_b["width"] / 2 - target_shape.Width / 2
            elif horizontal == "right":
                target_shape.Left = ref_b["left"] + ref_b["width"] - target_shape.Width

            if vertical == "top":
                target_shape.Top = ref_b["top"]
            elif vertical == "middle":
                target_shape.Top = ref_b["top"] + ref_b["height"] / 2 - target_shape.Height / 2
            elif vertical == "bottom":
                target_shape.Top = ref_b["top"] + ref_b["height"] - target_shape.Height

        try:
            # --- Mode: 1:1 paired (references is array) ---
            if references is not None:
                if len(references) != len(targets):
                    return {
                        "success": False,
                        "error": f"targets ({len(targets)}) and references ({len(references)}) must have same length",
                    }
                aligned = 0
                skipped = []
                for t_name, r_name in zip(targets, references):
                    t_shape = _find_shape(t_name)
                    if not t_shape:
                        skipped.append(t_name)
                        continue
                    r_shape = _find_shape(r_name)
                    if not r_shape:
                        return {"success": False, "error": f"Reference shape not found: {r_name}"}
                    _align_one(t_shape, _bounds(r_shape))
                    aligned += 1
                logger.info(f"[ActionExecutor] align_elements 1:1 paired, aligned={aligned}")
                return {"success": True, "aligned_count": aligned, "skipped": skipped}

            # --- Resolve target shapes (for other modes) ---
            target_shapes = []
            skipped = []
            for name in targets:
                s = _find_shape(name)
                if s:
                    target_shapes.append(s)
                else:
                    skipped.append(name)
            if not target_shapes:
                return {"success": False, "error": f"No target shapes found. Skipped: {skipped}"}

            # --- Mode: align to single reference ---
            if reference is not None:
                if reference == "slide":
                    ref_b = {
                        "left": 0.0, "top": 0.0,
                        "width": float(pres.PageSetup.SlideWidth),
                        "height": float(pres.PageSetup.SlideHeight),
                    }
                else:
                    r_shape = _find_shape(reference)
                    if not r_shape:
                        return {"success": False, "error": f"Reference shape not found: {reference}"}
                    ref_b = _bounds(r_shape)
                for ts in target_shapes:
                    _align_one(ts, ref_b)
                logger.info(
                    f"[ActionExecutor] align_elements to ref={reference}, "
                    f"h={horizontal}, v={vertical}, count={len(target_shapes)}"
                )
                return {"success": True, "aligned_count": len(target_shapes), "skipped": skipped}

            # --- Mode: mutual alignment (no reference) ---
            all_b = [_bounds(s) for s in target_shapes]

            if horizontal == "left":
                baseline = min(b["left"] for b in all_b)
                for s in target_shapes:
                    s.Left = baseline
            elif horizontal == "center":
                baseline = sum(b["left"] + b["width"] / 2 for b in all_b) / len(all_b)
                for s in target_shapes:
                    s.Left = baseline - s.Width / 2
            elif horizontal == "right":
                baseline = max(b["left"] + b["width"] for b in all_b)
                for s in target_shapes:
                    s.Left = baseline - s.Width

            if vertical == "top":
                baseline = min(b["top"] for b in all_b)
                for s in target_shapes:
                    s.Top = baseline
            elif vertical == "middle":
                baseline = sum(b["top"] + b["height"] / 2 for b in all_b) / len(all_b)
                for s in target_shapes:
                    s.Top = baseline - s.Height / 2
            elif vertical == "bottom":
                baseline = max(b["top"] + b["height"] for b in all_b)
                for s in target_shapes:
                    s.Top = baseline - s.Height

            logger.info(
                f"[ActionExecutor] align_elements mutual, "
                f"h={horizontal}, v={vertical}, count={len(target_shapes)}"
            )
            return {"success": True, "aligned_count": len(target_shapes), "skipped": skipped}

        except Exception as e:
            return {"success": False, "error": str(e)}

    def _handle_distribute_elements(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        Distribute multiple elements evenly.

        Params:
            slide: int|"current"
            targets: list[str]      - shapes to distribute (>= 3 recommended)
            handle_ids: list[str]   - alias for targets (backward compat)
            direction: str          - "horizontal" or "vertical"
            relative_to: str        - "slide" or "each_other" (default)
        """
        from .constants import DISTRIBUTE_TYPES

        slide_ref = params.get("slide", "current")
        handle_ids = params.get("targets") or params.get("handle_ids", [])
        direction = params.get("direction", "horizontal")
        relative_to = params.get("relative_to", "each_other")

        if len(handle_ids) < 2:
            return {"success": False, "error": "Need at least 2 handle_ids to distribute"}

        dist_cmd = DISTRIBUTE_TYPES.get(direction)
        if dist_cmd is None:
            return {
                "success": False,
                "error": f"Unknown direction: {direction}. Available: {list(DISTRIBUTE_TYPES.keys())}",
            }

        slide = resolve_slide_with_app(app, pres, slide_ref)

        indices = []
        missing = []
        for hid in handle_ids:
            found = False
            for i in range(1, slide.Shapes.Count + 1):
                try:
                    if slide.Shapes(i).Name == hid:
                        indices.append(i)
                        found = True
                        break
                except Exception:
                    continue
            if not found:
                missing.append(hid)

        if missing:
            return {"success": False, "error": f"Shapes not found: {missing}"}

        try:
            mso_relative = -1 if relative_to == "slide" else 0
            shape_range = slide.Shapes.Range(indices)
            shape_range.Distribute(dist_cmd, mso_relative)
            logger.info(
                f"[ActionExecutor] Distributed {len(handle_ids)} elements: "
                f"{direction}, relative_to={relative_to}"
            )
            return {"success": True, "distributed": handle_ids, "direction": direction}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ==================== Z-Order / Reorder ====================

    def _handle_reorder_elements(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        Change the z-order (layer stacking) of shapes.

        Three modes (by parameter priority):
          1. order[]               → batch reorder (bottom-to-top)
          2. command=above/below + reference → place relative to another shape
          3. command=bring_to_front/send_to_back/bring_forward/send_backward → basic

        Params:
            slide: int|"current"
            handle_id: str          - target shape (modes 2 & 3)
            command: str            - see above
            reference: str          - reference shape name (mode 2)
            order: list[str]        - shape names bottom-to-top (mode 1)
        """
        from .constants import Z_ORDER_COMMANDS

        slide_ref = params.get("slide", "current")
        slide = resolve_slide_with_app(app, pres, slide_ref)

        def _find(name):
            for i in range(1, slide.Shapes.Count + 1):
                try:
                    if slide.Shapes(i).Name == name:
                        return slide.Shapes(i)
                except Exception:
                    continue
            return None

        order = params.get("order")

        # --- Mode 1: batch reorder ---
        if order is not None:
            if len(order) < 2:
                return {"success": False, "error": "order needs at least 2 shape names"}

            shapes_to_reorder = []
            missing = []
            for name in order:
                s = _find(name)
                if s:
                    shapes_to_reorder.append((name, s))
                else:
                    missing.append(name)

            if len(shapes_to_reorder) < 2:
                return {"success": False, "error": f"Need at least 2 found shapes. Missing: {missing}"}

            try:
                for _name, s in shapes_to_reorder:
                    s.ZOrder(0)  # msoBringToFront
                logger.info(f"[ActionExecutor] reorder_elements batch, order={order}")
                return {"success": True, "reordered": order, "skipped": missing}
            except Exception as e:
                return {"success": False, "error": str(e)}

        # --- Modes 2 & 3: single shape ---
        handle_id = params.get("handle_id")
        command = params.get("command")

        if not handle_id:
            return {"success": False, "error": "handle_id is required (or use order[] for batch)"}
        if not command:
            return {"success": False, "error": "command is required"}

        target = _find(handle_id)
        if not target:
            return {"success": False, "error": f"Shape not found: {handle_id}"}

        try:
            # --- Mode 2: relative (above / below) ---
            if command in ("above", "below"):
                ref_name = params.get("reference")
                if not ref_name:
                    return {"success": False, "error": "reference is required for above/below"}
                ref_shape = _find(ref_name)
                if not ref_shape:
                    return {"success": False, "error": f"Reference shape not found: {ref_name}"}

                ref_pos = ref_shape.ZOrderPosition

                if command == "above":
                    target.ZOrder(0)  # msoBringToFront
                    total = slide.Shapes.Count
                    cur_pos = target.ZOrderPosition
                    steps = cur_pos - (ref_pos + 1)
                    if ref_pos < cur_pos:
                        ref_pos = ref_shape.ZOrderPosition
                        cur_pos = target.ZOrderPosition
                        steps = cur_pos - (ref_pos + 1)
                    for _ in range(max(0, steps)):
                        target.ZOrder(3)  # msoSendBackward
                else:  # below
                    target.ZOrder(1)  # msoSendToBack
                    cur_pos = target.ZOrderPosition
                    ref_pos = ref_shape.ZOrderPosition
                    steps = (ref_pos - 1) - cur_pos
                    for _ in range(max(0, steps)):
                        target.ZOrder(2)  # msoBringForward

                logger.info(
                    f"[ActionExecutor] reorder_elements {command} "
                    f"{handle_id} -> {ref_name}"
                )
                return {"success": True, "handle_id": handle_id, "command": command, "reference": ref_name}

            # --- Mode 3: basic command ---
            z_cmd = Z_ORDER_COMMANDS.get(command)
            if z_cmd is None:
                valid = list(Z_ORDER_COMMANDS.keys()) + ["above", "below"]
                return {"success": False, "error": f"Unknown command: {command}. Available: {valid}"}

            target.ZOrder(z_cmd)
            logger.info(f"[ActionExecutor] reorder_elements {command} {handle_id}")
            return {"success": True, "handle_id": handle_id, "command": command}

        except Exception as e:
            return {"success": False, "error": str(e)}

    # ==================== Animation / Transition ====================

    def _handle_add_shape_animation(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        Add an animation effect to a shape.

        Params:
            slide: int|"current"
            handle_id: str          - shape Name (or shape_index: int, 1-based)
            effect: str|int         - effect name or MsoAnimEffect integer
            category: str           - "entrance" (default) / "exit" / "emphasis"
            trigger: str            - "on_click" / "with_previous" / "after_previous"
            duration: float         - seconds (default 0.5)
            delay: float            - delay in seconds (default 0)
            direction: str          - optional, e.g. "from_left"
            text_unit: str          - optional, "paragraph" for per-paragraph build
            insert_at: int          - position in sequence (-1 = append, default)
        """
        from .constants import (
            ANIMATION_DIRECTIONS,
            ANIMATION_EFFECTS,
            ANIMATION_TRIGGERS,
        )

        slide_ref = params.get("slide", "current")
        handle_id = params.get("handle_id")
        shape_index = params.get("shape_index")
        effect = params.get("effect", "fade")
        category = params.get("category", "entrance")
        trigger = params.get("trigger", "on_click")
        duration = params.get("duration", 0.5)
        delay = params.get("delay", 0.0)
        direction = params.get("direction")
        text_unit = params.get("text_unit")
        insert_at = params.get("insert_at", -1)

        if not handle_id and not shape_index:
            return {"success": False, "error": "Missing 'handle_id' or 'shape_index'"}

        slide = resolve_slide_with_app(app, pres, slide_ref)

        # Find shape
        if handle_id:
            shape = find_shape_by_handle(slide, handle_id)
            if shape is None:
                return {"success": False, "error": f"Shape not found: '{handle_id}'"}
        else:
            idx = int(shape_index)
            if idx < 1 or idx > slide.Shapes.Count:
                return {"success": False, "error": f"shape_index {idx} out of range"}
            shape = slide.Shapes(idx)

        # Resolve effect ID
        if isinstance(effect, str):
            effect_id = ANIMATION_EFFECTS.get(effect)
            if effect_id is None:
                return {
                    "success": False,
                    "error": f"Unknown effect: {effect}. Available: {list(ANIMATION_EFFECTS.keys())}",
                }
        else:
            effect_id = int(effect)

        # Resolve trigger
        trigger_id = ANIMATION_TRIGGERS.get(trigger, 1)

        try:
            seq = slide.TimeLine.MainSequence
            anim = seq.AddEffect(shape, effect_id, 0, trigger_id)

            # Exit animation
            if category == "exit":
                anim.Exit = -1  # msoTrue

            # Timing
            anim.Timing.Duration = float(duration)
            if delay:
                anim.Timing.TriggerDelayTime = float(delay)

            # Direction
            if direction:
                dir_id = ANIMATION_DIRECTIONS.get(direction)
                if dir_id is not None:
                    try:
                        anim.EffectParameters.Direction = dir_id
                    except Exception:
                        pass

            # Text by paragraph
            if text_unit == "paragraph":
                try:
                    anim.EffectParameters.Amount = 0
                    # ppAnimateByParagraph: TextUnitEffect on the EffectInformation
                    # COM: use Effect.Behaviors to set text-level animation
                    # Simpler approach: ConvertToTextUnitEffect
                    anim_type_paragraph = 2  # ppAnimateByAllLevels
                    seq.ConvertToTextUnitEffect(anim, anim_type_paragraph)
                except Exception as e:
                    logger.debug(f"[ActionExecutor] text_unit=paragraph failed: {e}")

            # Reposition in sequence
            if insert_at > 0:
                try:
                    anim.MoveTo(insert_at)
                except Exception:
                    pass

            shape_name = handle_id or str(shape_index)
            logger.info(
                f"[ActionExecutor] Added animation '{effect}' ({category}) "
                f"to '{shape_name}', trigger={trigger}"
            )
            return {
                "success": True,
                "handle_id": shape_name,
                "effect": effect,
                "category": category,
                "trigger": trigger,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def _handle_clear_slide_animations(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        Clear all animations on a slide, or only those bound to a specific shape.

        Params:
            slide: int|"current"
            handle_id: str          - optional, clear only this shape's effects
        """
        slide_ref = params.get("slide", "current")
        handle_id = params.get("handle_id")

        slide = resolve_slide_with_app(app, pres, slide_ref)

        try:
            seq = slide.TimeLine.MainSequence
            removed = 0

            if handle_id:
                shape = find_shape_by_handle(slide, handle_id)
                if shape is None:
                    return {"success": False, "error": f"Shape not found: '{handle_id}'"}

                for i in range(seq.Count, 0, -1):
                    try:
                        if seq(i).Shape.Name == handle_id:
                            seq(i).Delete()
                            removed += 1
                    except Exception:
                        continue
            else:
                removed = seq.Count
                for i in range(seq.Count, 0, -1):
                    try:
                        seq(i).Delete()
                    except Exception:
                        continue

            logger.info(
                f"[ActionExecutor] Cleared {removed} animation(s) "
                f"on slide {slide.SlideIndex}"
                + (f" for '{handle_id}'" if handle_id else "")
            )
            return {"success": True, "removed": removed}
        except Exception as e:
            return {"success": False, "error": str(e)}

    def _handle_set_slide_transition(
        self, app, pres, params: Dict
    ) -> Dict[str, Any]:
        """
        Set the slide transition effect.

        Params:
            slide: int|"current"
            transition: str         - transition name (see SLIDE_TRANSITIONS)
            duration: float         - transition duration in seconds (default 0.5)
            advance_on_click: bool  - advance on mouse click (default True)
            advance_after_time: float - auto-advance after N seconds (0 = disabled)
        """
        from .constants import SLIDE_TRANSITIONS

        slide_ref = params.get("slide", "current")
        transition = params.get("transition", "fade")
        duration = params.get("duration", 0.5)
        advance_on_click = params.get("advance_on_click", True)
        advance_after_time = params.get("advance_after_time", 0)

        slide = resolve_slide_with_app(app, pres, slide_ref)

        entry_effect = SLIDE_TRANSITIONS.get(transition)
        if entry_effect is None:
            if isinstance(transition, int):
                entry_effect = transition
            else:
                return {
                    "success": False,
                    "error": f"Unknown transition: {transition}. Available: {list(SLIDE_TRANSITIONS.keys())}",
                }

        try:
            sst = slide.SlideShowTransition
            sst.EntryEffect = entry_effect
            sst.Duration = float(duration)
            sst.AdvanceOnClick = -1 if advance_on_click else 0

            if advance_after_time and float(advance_after_time) > 0:
                sst.AdvanceOnTime = -1  # msoTrue
                sst.AdvanceTime = float(advance_after_time)
            else:
                sst.AdvanceOnTime = 0  # msoFalse

            logger.info(
                f"[ActionExecutor] Set transition '{transition}' on slide {slide.SlideIndex}, "
                f"duration={duration}s"
            )
            return {"success": True, "transition": transition, "duration": duration}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ==================== Property Application ====================

    def _apply_properties(self, shape, props: Dict) -> None:
        """Apply a dict of properties to a shape.

        Supported keys
        ──────────────
        Geometry:  x, y, width, height, rotation, visible, z_order
        Fill:      fill_color  (str, mutually exclusive with fill_gradient)
                   fill_gradient (dict: type, angle, stops[])
        Line:      line_color (str, mutually exclusive with line_gradient)
                   line_gradient (dict: same schema as fill_gradient)
                   line_weight (float)
        Shadow:    shadow (dict or None to remove)
        Text:      text, font_name, font_size, font_bold, font_italic,
                   font_color, text_align
        """
        # ── Geometry ──────────────────────────────────────────────
        if "x" in props:
            shape.Left = float(props["x"])
        if "y" in props:
            shape.Top = float(props["y"])
        if "width" in props:
            shape.Width = float(props["width"])
        if "height" in props:
            shape.Height = float(props["height"])
        if "rotation" in props:
            shape.Rotation = float(props["rotation"])
        if "visible" in props:
            shape.Visible = bool(props["visible"])

        if "z_order" in props:
            cmd_name = props["z_order"]
            cmd = Z_ORDER_COMMANDS.get(cmd_name)
            if cmd is not None:
                shape.ZOrder(cmd)
            else:
                raise ValueError(
                    f"Unknown z_order command: '{cmd_name}'. "
                    f"Available: {list(Z_ORDER_COMMANDS.keys())}"
                )

        # ── Fill (solid OR gradient, mutually exclusive) ──────────
        if "fill_gradient" in props:
            self._apply_fill_gradient(shape, props["fill_gradient"])
        elif "fill_color" in props:
            color = parse_color(props["fill_color"])
            if color is not None:
                shape.Fill.Solid()
                shape.Fill.ForeColor.RGB = color
            else:
                shape.Fill.Background()

        # ── Line (solid color OR gradient, mutually exclusive) ─────
        if "line_gradient" in props:
            self._apply_line_gradient(shape, props["line_gradient"])
        elif "line_color" in props:
            color = parse_color(props["line_color"])
            if color is not None:
                shape.Line.Visible = True
                shape.Line.ForeColor.RGB = color
            else:
                shape.Line.Visible = False
        if "line_weight" in props:
            shape.Line.Weight = float(props["line_weight"])

        # ── Shadow ────────────────────────────────────────────────
        if "shadow" in props:
            self._apply_shadow(shape, props["shadow"])

        # ── Text ──────────────────────────────────────────────────
        # rich_text takes precedence over plain text/font props
        if "rich_text" in props and shape.HasTextFrame:
            self._apply_rich_text(shape, props["rich_text"])
        else:
            if "text" in props and shape.HasTextFrame:
                shape.TextFrame.TextRange.Text = str(props["text"])

            if shape.HasTextFrame:
                font = shape.TextFrame.TextRange.Font
                if "font_name" in props:
                    font.Name = str(props["font_name"])
                if "font_size" in props:
                    font.Size = float(props["font_size"])
                if "font_bold" in props:
                    font.Bold = bool(props["font_bold"])
                if "font_italic" in props:
                    font.Italic = bool(props["font_italic"])
                if "font_color" in props:
                    color = parse_color(props["font_color"])
                    if color is not None:
                        font.Color.RGB = color

        # text_formats: apply formatting to ranges of EXISTING text
        if "text_formats" in props and shape.HasTextFrame:
            self._apply_text_formats(shape, props["text_formats"])

        if "text_align" in props and shape.HasTextFrame:
            align = TEXT_ALIGN.get(props["text_align"])
            if align is not None:
                for i in range(1, shape.TextFrame.TextRange.Paragraphs().Count + 1):
                    shape.TextFrame.TextRange.Paragraphs(i).ParagraphFormat.Alignment = align

        # ── Table ──────────────────────────────────────────────────
        if shape.HasTable:
            table = shape.Table
            if "table_style" in props:
                style_guid = resolve_table_style(props["table_style"])
                if style_guid:
                    table.ApplyStyle(style_guid)

            if "first_row_header" in props:
                table.FirstRow = bool(props["first_row_header"])
            if "first_col" in props:
                table.FirstCol = bool(props["first_col"])
            if "banded_rows" in props:
                table.HorizBanding = bool(props["banded_rows"])
            if "banded_cols" in props:
                table.VertBanding = bool(props["banded_cols"])

            if "cell_format" in props:
                from .macro_executor import MacroExecutor
                num_rows = table.Rows.Count
                num_cols = table.Columns.Count
                MacroExecutor._apply_cell_format(
                    table, props["cell_format"], num_rows, num_cols
                )

    # ── rich_text / text_formats helpers ─────────────────────────

    @staticmethod
    def _apply_font_props_to_range(chars, seg: Dict) -> None:
        """Apply font properties from *seg* dict to a COM TextRange *chars*."""
        font = chars.Font
        if "font_name" in seg:
            font.Name = str(seg["font_name"])
        if "font_size" in seg:
            font.Size = float(seg["font_size"])
        if "font_bold" in seg:
            font.Bold = bool(seg["font_bold"])
        if "font_italic" in seg:
            font.Italic = bool(seg["font_italic"])
        if "font_underline" in seg:
            font.Underline = bool(seg["font_underline"])
        if "font_color" in seg:
            color = parse_color(seg["font_color"])
            if color is not None:
                font.Color.RGB = color

    def _apply_rich_text(self, shape, segments: List[Dict]) -> None:
        """Replace text with per-segment formatting.

        Each segment: ``{"text": "...", "font_name": ..., "font_size": ...,
        "font_bold": ..., "font_italic": ..., "font_underline": ...,
        "font_color": "#RRGGBB", "fill_gradient": {...}, "highlight": "#RRGGBB"}``

        Only ``text`` is required; all formatting keys are optional.
        Segments are concatenated to form the full text, then each segment's
        formatting is applied via ``Characters(start, length)``.
        """
        if not segments:
            return

        full_text = "".join(seg.get("text", "") for seg in segments)
        shape.TextFrame.TextRange.Text = full_text

        tr = shape.TextFrame.TextRange
        offset = 1  # COM Characters() is 1-based
        gradient_segments = []

        for seg in segments:
            text = seg.get("text", "")
            length = len(text)
            if length == 0:
                continue
            chars = tr.Characters(offset, length)
            self._apply_font_props_to_range(chars, seg)

            if "fill_gradient" in seg:
                gradient_segments.append((offset, length, seg["fill_gradient"]))

            offset += length

        # Apply text gradients using the whole-range-then-override approach:
        # TextFrame2.TextRange.Characters() sub-ranges often fail on
        # late-bound COM, so we apply gradient to the FULL range and then
        # override non-gradient segments back to solid color.
        if gradient_segments:
            self._apply_segmented_text_gradient(
                shape, segments, gradient_segments, full_text
            )

    def _apply_text_formats(self, shape, formats: List[Dict]) -> None:
        """Apply formatting to ranges of existing text without replacing it.

        Each entry can locate text by index or by substring match::

            {"start": 0, "length": 5, "font_bold": true, "font_color": "#FF0000"}
            {"match": "keyword", "font_bold": true, "font_color": "#FF0000"}
            {"match": "keyword", "nth": 2, ...}   # 2nd occurrence (1-based, default 1)

        ``start`` is 0-based; it is converted to 1-based for COM internally.
        ``match`` does a substring search on the current text; if both ``start``
        and ``match`` are given, ``match`` takes precedence.
        """
        if not formats:
            return

        current_text = shape.TextFrame.TextRange.Text
        tr = shape.TextFrame.TextRange

        for fmt in formats:
            start_0 = None
            length = None

            if "match" in fmt:
                target = str(fmt["match"])
                nth = int(fmt.get("nth", 1))
                idx = -1
                for _ in range(nth):
                    idx = current_text.find(target, idx + 1)
                    if idx == -1:
                        break
                if idx == -1:
                    logger.warning("text_formats match %r (nth=%d) not found", target, nth)
                    continue
                start_0 = idx
                length = len(target)
            elif "start" in fmt:
                start_0 = int(fmt["start"])
                length = int(fmt.get("length", len(current_text) - start_0))
            else:
                continue

            if start_0 < 0 or start_0 >= len(current_text):
                continue
            length = min(length, len(current_text) - start_0)
            if length <= 0:
                continue

            chars = tr.Characters(start_0 + 1, length)  # COM is 1-based
            self._apply_font_props_to_range(chars, fmt)

            if "text" in fmt:
                chars.Text = str(fmt["text"])

            if "fill_gradient" in fmt:
                try:
                    tf2_chars = shape.TextFrame2.TextRange.Characters(start_0 + 1, length)
                    self._apply_text_range_gradient(tf2_chars, fmt["fill_gradient"])
                except Exception as e:
                    logger.warning(
                        "text_formats fill_gradient on sub-range failed: %s", e
                    )

    def _apply_segmented_text_gradient(
        self, shape, segments: List[Dict],
        gradient_segments: list, full_text: str,
    ) -> None:
        """Apply per-run gradient/solid fills via python-pptx donor + COM paste.

        ``TextFrame2.TextRange.Characters()`` does not work on late-bound COM,
        so per-character gradient cannot be set through COM APIs alone.
        Instead we:

        1. Build a donor ``.pptx`` with a text shape whose runs have the
           exact ``<a:gradFill>`` / ``<a:solidFill>`` XML in each ``<a:rPr>``.
        2. Open the donor hidden in COM, copy the formatted text range.
        3. Paste into the target shape (preserves per-run fill formatting).
        4. Re-apply non-fill font properties (bold/size/etc.) via COM since
           paste may alter them.
        """
        import time
        from pptx import Presentation as PptxPresentation
        from pptx.util import Inches
        from pptx.oxml.ns import qn
        from lxml import etree

        # ── Build donor .pptx with per-run XML fills ─────────────
        prs = PptxPresentation()
        pptx_slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = pptx_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        tf = txBox.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.clear()

        for i, seg in enumerate(segments):
            text = seg.get("text", "")
            if not text:
                continue
            run = p.add_run() if (i > 0 or not p.runs) else p.runs[0]
            run.text = text

            rPr = run._r.get_or_add_rPr()
            if seg.get("font_bold"):
                rPr.set("b", "1")
            if seg.get("font_italic"):
                rPr.set("i", "1")
            if seg.get("font_underline"):
                rPr.set("u", "sng")
            sz = seg.get("font_size")
            if sz:
                rPr.set("sz", str(int(float(sz) * 100)))
            fname = seg.get("font_name")
            if fname:
                latin = etree.SubElement(rPr, qn("a:latin"))
                latin.set("typeface", fname)

            for old in rPr.findall(qn("a:solidFill")) + rPr.findall(qn("a:gradFill")):
                rPr.remove(old)

            if "fill_gradient" in seg:
                gdef = seg["fill_gradient"]
                stops = gdef.get("stops", [])
                gradFill = etree.SubElement(rPr, qn("a:gradFill"))
                gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
                for stop in sorted(stops, key=lambda s: s["position"]):
                    gs = etree.SubElement(gsLst, qn("a:gs"))
                    gs.set("pos", str(int(float(stop["position"]) * 100000)))
                    srgb = etree.SubElement(gs, qn("a:srgbClr"))
                    srgb.set("val", stop["color"].lstrip("#"))
                    opacity = stop.get("opacity", 1.0)
                    if opacity < 1.0:
                        alpha = etree.SubElement(srgb, qn("a:alpha"))
                        alpha.set("val", str(int(opacity * 100000)))
                grad_type = gdef.get("type", "linear")
                if grad_type == "radial":
                    path_el = etree.SubElement(gradFill, qn("a:path"))
                    path_el.set("path", "circle")
                    ftr = etree.SubElement(path_el, qn("a:fillToRect"))
                    for attr in ("l", "t", "r", "b"):
                        ftr.set(attr, "50000")
                else:
                    lin = etree.SubElement(gradFill, qn("a:lin"))
                    lin.set("ang", str(int(float(gdef.get("angle", 0)) * 60000)))
                    lin.set("scaled", "1")
            elif "font_color" in seg:
                solidFill = etree.SubElement(rPr, qn("a:solidFill"))
                srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
                srgb.set("val", seg["font_color"].lstrip("#"))

        donor_path = os.path.join(
            tempfile.gettempdir(), f"_ppt_rich_text_donor_{os.getpid()}.pptx"
        )
        prs.save(donor_path)

        # ── COM: open donor, copy formatted text, paste into target ──
        logger.info("[rich_text] Using donor pptx + paste approach for per-run gradient")
        import win32com.client
        app = win32com.client.GetActiveObject("PowerPoint.Application")
        donor_pres = None
        try:
            donor_pres = app.Presentations.Open(
                donor_path, ReadOnly=True, WithWindow=False
            )
            donor_shape = donor_pres.Slides(1).Shapes(1)
            donor_text = donor_shape.TextFrame.TextRange.Text
            logger.info("[rich_text] Donor text: %r", donor_text[:60])

            donor_shape.TextFrame.TextRange.Copy()
            time.sleep(0.3)

            # Deselect any current selection first to avoid stale state
            try:
                app.ActiveWindow.Selection.Unselect()
            except Exception:
                pass
            time.sleep(0.1)

            shape.TextFrame.TextRange.Text = ""
            shape.Select()
            time.sleep(0.3)

            # Enter text editing mode by double-clicking into the text frame
            try:
                shape.TextFrame.TextRange.Select()
            except Exception:
                pass
            time.sleep(0.3)

            app.ActiveWindow.View.Paste()
            time.sleep(0.2)

            # Click away to deselect, preventing stale clipboard state
            try:
                app.ActiveWindow.Selection.Unselect()
            except Exception:
                pass

            result_text = shape.TextFrame.TextRange.Text
            logger.info("[rich_text] Result text: %r", result_text[:60])
        except Exception as e:
            logger.error("[rich_text] Donor paste failed: %s", e)
            raise
        finally:
            if donor_pres:
                try:
                    donor_pres.Close()
                except Exception:
                    pass
            try:
                os.unlink(donor_path)
            except Exception:
                pass

    @staticmethod
    def _apply_text_range_gradient(tf2_range, gradient: Dict) -> None:
        """Apply gradient fill to a TextFrame2 text range (Characters/full range)."""
        stops = gradient.get("stops", [])
        if len(stops) < 2:
            return
        stops = sorted(stops, key=lambda s: s["position"])
        grad_type = gradient.get("type", "linear")

        ff = tf2_range.Font.Fill
        ff.Visible = True
        if grad_type == "radial":
            ff.TwoColorGradient(MSO_GRADIENT_FROM_CENTER, 1)
        else:
            ff.TwoColorGradient(MSO_GRADIENT_HORIZONTAL, 1)

        gs = ff.GradientStops
        c0 = parse_color(stops[0]["color"]) or 0
        cN = parse_color(stops[-1]["color"]) or 0
        gs(1).Color.RGB = c0
        gs(1).Position = float(stops[0]["position"])
        gs(1).Transparency = 1.0 - float(stops[0].get("opacity", 1.0))
        gs(2).Color.RGB = cN
        gs(2).Position = float(stops[-1]["position"])
        gs(2).Transparency = 1.0 - float(stops[-1].get("opacity", 1.0))

        for stop in stops[1:-1]:
            c = parse_color(stop["color"]) or 0
            pos = float(stop["position"])
            gs.Insert(c, pos)
            for idx in range(1, gs.Count + 1):
                if abs(gs(idx).Position - pos) < 0.005:
                    gs(idx).Transparency = 1.0 - float(stop.get("opacity", 1.0))
                    break

        if grad_type != "radial":
            angle = float(gradient.get("angle", 0))
            ff.GradientAngle = angle % 360

    # ── fill_gradient helper ──────────────────────────────────────

    @staticmethod
    def _apply_fill_gradient(shape, gradient: Dict) -> None:
        """Apply a gradient fill to *shape*.

        ``gradient`` schema::

            {
                "type": "linear" | "radial",
                "angle": float,          # degrees, only for linear
                "stops": [
                    {"position": 0.0, "color": "#RRGGBB", "opacity": 1.0},
                    {"position": 1.0, "color": "#RRGGBB", "opacity": 1.0},
                    ...
                ]
            }
        """
        stops = gradient.get("stops", [])
        if len(stops) < 2:
            raise ValueError("fill_gradient requires at least 2 stops")

        stops = sorted(stops, key=lambda s: s["position"])
        grad_type = gradient.get("type", "linear")

        fill = shape.Fill

        if grad_type == "radial":
            fill.TwoColorGradient(MSO_GRADIENT_FROM_CENTER, 1)
        else:
            fill.TwoColorGradient(MSO_GRADIENT_HORIZONTAL, 1)

        gs = fill.GradientStops

        # Set the two default stops first
        c0 = parse_color(stops[0]["color"]) or 0
        cN = parse_color(stops[-1]["color"]) or 0
        gs(1).Color.RGB = c0
        gs(1).Position = float(stops[0]["position"])
        gs(1).Transparency = 1.0 - float(stops[0].get("opacity", 1.0))
        gs(2).Color.RGB = cN
        gs(2).Position = float(stops[-1]["position"])
        gs(2).Transparency = 1.0 - float(stops[-1].get("opacity", 1.0))

        # Insert middle stops
        for stop in stops[1:-1]:
            c = parse_color(stop["color"]) or 0
            pos = float(stop["position"])
            gs.Insert(c, pos)
            # Set transparency on the newly inserted stop
            for idx in range(1, gs.Count + 1):
                if abs(gs(idx).Position - pos) < 0.005:
                    gs(idx).Transparency = 1.0 - float(stop.get("opacity", 1.0))
                    break

        # Set angle for linear gradients
        if grad_type != "radial":
            angle = float(gradient.get("angle", 0))
            fill.GradientAngle = angle % 360

    # ── shadow helper ─────────────────────────────────────────────

    @staticmethod
    def _apply_shadow(shape, shadow) -> None:
        """Apply or remove a shadow effect.

        If *shadow* is ``None``, the shadow is hidden.
        Otherwise it should be a dict with optional keys:
        ``color``, ``blur``, ``offset_x``, ``offset_y``, ``opacity``.
        """
        if shadow is None:
            shape.Shadow.Visible = False
            return

        shd = shape.Shadow
        shd.Visible = True

        if "color" in shadow:
            color = parse_color(shadow["color"])
            if color is not None:
                shd.ForeColor.RGB = color

        if "blur" in shadow:
            shd.Blur = float(shadow["blur"])

        if "offset_x" in shadow:
            shd.OffsetX = float(shadow["offset_x"])

        if "offset_y" in shadow:
            shd.OffsetY = float(shadow["offset_y"])

        if "opacity" in shadow:
            shd.Transparency = 1.0 - float(shadow["opacity"])

    # ── line_gradient helper (python-pptx XML + COM PickUp/Apply) ─

    @staticmethod
    def _apply_line_gradient(shape, gradient: Dict) -> None:
        """Apply a gradient stroke to *shape* via the donor-shape technique.

        COM's ``LineFormat`` does not expose a ``Fill`` property, so we
        create a temporary ``.pptx`` with a donor shape whose ``<a:ln>``
        contains ``<a:gradFill>`` (valid OOXML), open it hidden via COM,
        then ``PickUp()`` → ``Apply()`` to transfer the format in real-time.

        ``gradient`` uses the same schema as ``fill_gradient``::

            {"type": "linear", "angle": 45,
             "stops": [{"position": 0, "color": "#FF4D00", "opacity": 1}, ...]}
        """
        from pptx import Presentation as PptxPresentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches
        from pptx.oxml.ns import qn
        from lxml import etree
        import win32com.client

        stops = gradient.get("stops", [])
        if len(stops) < 2:
            raise ValueError("line_gradient requires at least 2 stops")
        stops = sorted(stops, key=lambda s: s["position"])

        grad_type = gradient.get("type", "linear")
        weight_pt = gradient.get("weight")

        # ── Build donor .pptx ────────────────────────────────
        pptx_pres = PptxPresentation()
        pptx_slide = pptx_pres.slides.add_slide(pptx_pres.slide_layouts[6])
        donor = pptx_slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1), Inches(1), Inches(2), Inches(2),
        )
        donor.fill.background()

        sp = donor._element
        spPr = None
        for child in sp:
            if "spPr" in child.tag:
                spPr = child
                break

        ln = spPr.find(qn("a:ln"))
        if ln is None:
            ln = etree.SubElement(spPr, qn("a:ln"))

        w_emu = int(float(weight_pt) * 12700) if weight_pt else int(shape.Line.Weight * 12700)
        ln.set("w", str(w_emu))

        for child in list(ln):
            if child.tag in [qn("a:solidFill"), qn("a:noFill"), qn("a:gradFill")]:
                ln.remove(child)

        gradFill = etree.SubElement(ln, qn("a:gradFill"))
        gsLst = etree.SubElement(gradFill, qn("a:gsLst"))

        for stop in stops:
            gs = etree.SubElement(gsLst, qn("a:gs"))
            gs.set("pos", str(int(float(stop["position"]) * 100000)))
            srgb = etree.SubElement(gs, qn("a:srgbClr"))
            color_hex = stop["color"].lstrip("#")
            srgb.set("val", color_hex)
            opacity = stop.get("opacity", 1.0)
            if opacity < 1.0:
                alpha = etree.SubElement(srgb, qn("a:alpha"))
                alpha.set("val", str(int(opacity * 100000)))

        if grad_type == "radial":
            path = etree.SubElement(gradFill, qn("a:path"))
            path.set("path", "circle")
            fillToRect = etree.SubElement(path, qn("a:fillToRect"))
            fillToRect.set("l", "50000")
            fillToRect.set("t", "50000")
            fillToRect.set("r", "50000")
            fillToRect.set("b", "50000")
        else:
            angle = float(gradient.get("angle", 0))
            lin_elem = etree.SubElement(gradFill, qn("a:lin"))
            lin_elem.set("ang", str(int(angle * 60000)))
            lin_elem.set("scaled", "1")

        # Save to temp file
        donor_path = os.path.join(
            tempfile.gettempdir(), f"_ppt_line_grad_donor_{os.getpid()}.pptx"
        )
        pptx_pres.save(donor_path)

        # ── Capture current fill & shadow state before PickUp/Apply ──
        # PickUp/Apply transfers ALL formatting (fill + line + shadow),
        # so we must snapshot and restore the properties we don't want to change.
        saved_fill = PPTActionExecutor._capture_fill_state(shape)
        saved_shadow = PPTActionExecutor._capture_shadow_state(shape)

        # ── COM PickUp/Apply ─────────────────────────────────
        app = win32com.client.GetActiveObject("PowerPoint.Application")
        donor_pres = None
        try:
            donor_pres = app.Presentations.Open(donor_path, ReadOnly=True, WithWindow=False)
            donor_shape = donor_pres.Slides(1).Shapes(1)
            donor_shape.PickUp()
            shape.Apply()
        finally:
            if donor_pres:
                try:
                    donor_pres.Close()
                except Exception:
                    pass
            try:
                os.unlink(donor_path)
            except Exception:
                pass

        # ── Restore fill & shadow that PickUp/Apply overwrote ──
        PPTActionExecutor._restore_fill_state(shape, saved_fill)
        PPTActionExecutor._restore_shadow_state(shape, saved_shadow)

    # ── fill / shadow snapshot helpers for PickUp/Apply isolation ──

    @staticmethod
    def _capture_fill_state(shape) -> Dict:
        """Snapshot a shape's fill so it can be restored after PickUp/Apply."""
        try:
            fill = shape.Fill
            fill_type = fill.Type
            # msoFillBackground = 5 (no fill / background)
            if fill_type == 5:
                return {"type": "background"}
            # msoFillSolid = 1
            if fill_type == 1:
                return {
                    "type": "solid",
                    "rgb": fill.ForeColor.RGB,
                    "transparency": fill.Transparency,
                }
            # msoFillGradient = 3
            if fill_type == 3:
                gs = fill.GradientStops
                stops = []
                for i in range(1, gs.Count + 1):
                    stops.append({
                        "position": gs(i).Position,
                        "rgb": gs(i).Color.RGB,
                        "transparency": gs(i).Transparency,
                    })
                try:
                    angle = fill.GradientAngle
                except Exception:
                    angle = 0
                return {
                    "type": "gradient",
                    "gradient_style": fill.GradientStyle,
                    "gradient_variant": 1,
                    "stops": stops,
                    "angle": angle,
                }
            # msoFillPatterned = 2, msoFillTextured = 4, msoFillPicture = 6
            return {"type": "unknown", "fill_type": fill_type}
        except Exception:
            return {"type": "unknown"}

    @staticmethod
    def _restore_fill_state(shape, state: Dict) -> None:
        """Restore a shape's fill from a snapshot captured by _capture_fill_state."""
        fill_type = state.get("type", "unknown")
        if fill_type == "unknown":
            return
        try:
            fill = shape.Fill
            if fill_type == "background":
                fill.Background()
                return
            if fill_type == "solid":
                fill.Solid()
                fill.ForeColor.RGB = state["rgb"]
                fill.Transparency = state.get("transparency", 0)
                return
            if fill_type == "gradient":
                style = state.get("gradient_style", MSO_GRADIENT_HORIZONTAL)
                variant = state.get("gradient_variant", 1)
                fill.TwoColorGradient(style, variant)
                gs = fill.GradientStops
                saved_stops = state.get("stops", [])
                # Restore existing stops
                for i, ss in enumerate(saved_stops[:gs.Count], 1):
                    gs(i).Color.RGB = ss["rgb"]
                    gs(i).Position = ss["position"]
                    gs(i).Transparency = ss.get("transparency", 0)
                # Insert additional stops if there were more
                for ss in saved_stops[gs.Count:]:
                    gs.Insert(ss["rgb"], ss["position"])
                    for j in range(1, gs.Count + 1):
                        if abs(gs(j).Position - ss["position"]) < 0.005:
                            gs(j).Transparency = ss.get("transparency", 0)
                            break
                try:
                    fill.GradientAngle = state.get("angle", 0)
                except Exception:
                    pass
        except Exception:
            pass

    @staticmethod
    def _capture_shadow_state(shape) -> Dict:
        """Snapshot a shape's shadow properties."""
        try:
            shd = shape.Shadow
            return {
                "visible": shd.Visible,
                "fore_rgb": shd.ForeColor.RGB if shd.Visible else 0,
                "blur": shd.Blur if shd.Visible else 0,
                "offset_x": shd.OffsetX if shd.Visible else 0,
                "offset_y": shd.OffsetY if shd.Visible else 0,
                "transparency": shd.Transparency if shd.Visible else 0,
            }
        except Exception:
            return {}

    @staticmethod
    def _restore_shadow_state(shape, state: Dict) -> None:
        """Restore a shape's shadow from a snapshot."""
        if not state:
            return
        try:
            shd = shape.Shadow
            if not state.get("visible", False):
                shd.Visible = False
                return
            shd.Visible = True
            shd.ForeColor.RGB = state.get("fore_rgb", 0)
            shd.Blur = state.get("blur", 0)
            shd.OffsetX = state.get("offset_x", 0)
            shd.OffsetY = state.get("offset_y", 0)
            shd.Transparency = state.get("transparency", 0)
        except Exception:
            pass
